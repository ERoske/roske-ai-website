"""Generate PowerPoint presentations from structured data."""
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_presentation(data, output_path):
    """Create a PowerPoint presentation from structured data.

    Args:
        data: Dictionary with presentation structure
        output_path: Path to save the .pptx file
    """
    prs = Presentation()
    # Set to 16:9 widescreen format
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Set default font
    for i, slide_data in enumerate(data['slides']):
        # Use blank layout for Beyond Bullet Points style
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Add headline text box (white text, large font)
        left = Inches(0.5)
        top = Inches(3)
        width = Inches(9)
        height = Inches(1.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = slide_data['headline']
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Add semi-transparent background to text for readability
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)
        textbox.fill.transparency = 0.3

        # Add speaker notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_data['speaker_notes']

        # Add visual description as a note at the top
        visual_note = f"[IMAGE: {slide_data['visual']}]\n\n"
        text_frame.text = visual_note + text_frame.text

    # Save presentation
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python generate_pptx.py <json_file>")
        sys.exit(1)

    json_file = sys.argv[1]

    # Load presentation data
    with open(json_file, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # Generate output filename in Presentations subdirectory
    # Remove invalid Windows filename characters
    title_safe = data['title'].replace(' ', '_').replace(':', '-').replace('?', '').replace('/', '-').replace('\\', '-').replace('*', '').replace('"', '').replace('<', '').replace('>', '').replace('|', '')[:50]
    base_dir = Path(__file__).parent.parent.parent.parent  # Navigate to personal_assistant/
    output_path = base_dir / 'outputs' / 'Presentations' / f"{title_safe}.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Create presentation
    create_presentation(data, str(output_path))

    print(f"[OK] PowerPoint created: {output_path}")
    print(f"  Title: {data['title']}")
    print(f"  Slides: {len(data['slides'])}")
    print(f"\nNote: This presentation uses text-only slides.")
    print("To add images, open the file and insert images as backgrounds for each slide.")
    print("Visual descriptions are included in speaker notes for reference.")
